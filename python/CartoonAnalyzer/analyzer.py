#!/usr/bin/env python3
# -*- coding: UTF-8 -*-
#  analyzer.py
#  Copyright (C) 2019 Zhang Maiyun <me@maiyun.me>
#
#  This program is free software; you can redistribute it and/or modify
#  it under the terms of the GNU General Public License as published by
#  the Free Software Foundation; either version 3 of the License, or
#  (at your option) any later version.
#
#  This program is distributed in the hope that it will be useful,
#  but WITHOUT ANY WARRANTY; without even the implied warranty of
#  MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE.  See the
#  GNU General Public License for more details.
#
#  You should have received a copy of the GNU General Public License
#  along with this program.  If not, see <https://www.gnu.org/licenses/>.
#

"""Analyze YUSS 2020 cartoon specs and make instruction tables."""

import argparse
import json
import sys

import docx
import pptx
from docx.oxml import parse_xml
from docx.oxml.ns import nsdecls
from docx.shared import Cm


def rgbFromStr(s):
    # s starts with a #.
    r, g, b = int(s[1:3], 16), int(s[3:5], 16), int(s[5:7], 16)
    return r, g, b


def findNearestColorName(RGB_tuple, Map):
    R = RGB_tuple[0]
    G = RGB_tuple[1]
    B = RGB_tuple[2]
    mindiff = None
    for d in Map:
        r, g, b = rgbFromStr(Map[d])
        diff = abs(R - r)*256 + abs(G-g) * 256 + abs(B - b) * 256
        if mindiff is None or diff < mindiff:
            mindiff = diff
            mincolorname = d
    return mincolorname, Map[mincolorname]


def readppt(address, numbering):
    """Read a pptx cartoon specification.
    address: the filename to the pptx
    numbering: the numbering scheme, "m.8" or "int"
    Returns the data and the number of slides"""
    result = []
    ppt = pptx.Presentation(address)
    for i, slide in enumerate(ppt.slides):
        # The slide index: m.n where 0<n<=8 and m increases
        m = i // 8 + 1
        n = i % 8 + 1
        for shape in slide.shapes:
            if shape.shape_type == 19:
                table = shape.table
                break
        for nr, row in enumerate(table.rows):
            if nr == 0 or nr == 27:
                continue  # skip title bar
            for nc, cell in enumerate(row.cells):
                if nc == 0 or nc == 41:
                    continue
                colorMap = {
                    "W": "#FFFFFF",
                    "R": "#FF0000",
                    "B": "#00B0F0",
                    "O": "#FF6600",
                    "Y": "#FFFF00",
                    "K": "#7C7C7C"
                }
                # if there is annotation, use it
                color_anno = cell.text
                code_anno = colorMap[color_anno][1:] if color_anno else ''
                # fall back to grip reading, might fail
                try:
                    if cell.fill.fore_color.type == 1:
                        color, code = findNearestColorName(
                            cell.fill.fore_color.rgb, colorMap)
                    elif cell.fill.fore_color.type == 2:
                        # XXX: Read ppt/theme/theme1.xml
                        themeMap = {
                            13: "K",
                            14: "W"
                        }
                        try:
                            color = themeMap[cell.fill.fore_color.theme_color]
                        except KeyError as e:
                            print(
                                f"KeyError at {i}/{nr}/{nc}", file=sys.stderr)
                            raise
                        code = colorMap[color]
                    code = code[1:]
                    if color != color_anno:
                        choice = input(
                            f"The annotation at {i}/{nr}/{nc}, which is `{color_anno}' disagrees with the color `{color},{code}' analyzed by the program. Press 1 to proceed with the annotation and 2 to go with the program, with the annotation corrected, or 3 to temporarily use the program result: ")
                        if choice == '1':
                            color = color_anno
                            code = code_anno
                        elif choice == '2':
                            cell.text = color
                        elif choice == '3':
                            pass
                        else:
                            ppt.save(address)
                            raise Exception
                except AttributeError:
                    # print(i,nr,nc,cell.fill.fore_color.theme_color)
                    color, code = "unknown", "FFFFFF"
                if 0 < nc < 21:  # [1,20], read the left title
                    label = 0
                elif 20 < nc < 41:  # [21,40], read the right title
                    label = 41
                # The class number as a string
                classn = row.cells[label].text[:-1]
                # The row number, which is a alphabet
                coord = row.cells[label].text[-1]
                # add the column number
                coord += table.rows[0].cells[nc].text
                result.append({
                    "color": color,
                    "colorcode": code,
                    "class": classn,
                    "coord": coord,
                    "index": f"{m}.{n}" if numbering == "m.8" else str(i+1)
                })
    ppt.save(address)
    return result, len(ppt.slides)


def gendoc(data, length, numbering, outaddr, only, lyrics):
    """Generate a docx file containing the tables of instruction.
    data: the data from readppt()
    length: the number of slides
    numbering: the numbering scheme
    outaddr: output path
    only: list of the classes to be made
    lyrics: the lyrics file or None if not needed
    returns nothing"""
    doc = docx.Document()
    doc.core_properties.author = "YUSS 2020 Cartoon Analyzer"
    # 26 classes
    gen = only if only else range(1, 27)
    if lyrics:
        lyricslst = json.load(open(lyrics))
    print(f"Generating for classes: {gen}")
    for classn in gen:
        print(classn)
        for person in range(40):
            # Same lookup strategy
            coord = 'A' if person < 20 else 'B'
            coord += str(person % 20 + 1)
            # I merge two tables into one to save paper
            # There is a empty column between two
            if person % 2 == 0:
                offset = 0
                table = doc.add_table(rows=9, cols=9)
            else:
                offset = 4 + 1
            # Add grid lines
            table.style = 'Table Grid'
            r = table.rows[0].cells
            r[offset + 0].text = "班级"
            r[offset + 1].text = str(classn)
            r[offset + 2].text = "编码"
            r[offset + 3].text = coord
            r = table.rows[1].cells
            r[offset + 0].text = "姓名"
            r[offset + 1].merge(r[offset + 3])
            # Change if needed(XXX)
            # effectiverows = nrows - the number of class/name lines
            metalines = 2
            er = 9 - metalines
            for i in range(length):
                # Only for m.n numbering
                m = i // 8 + 1
                n = i % 8 + 1
                # Table index from 0, this is math
                xindex = offset + (i // 7) * 2
                xdata = offset + (i // 7) * 2 + 1
                y = metalines + i % 7
                index = f"{m}.{n}" if numbering == "m.8" else f"{i+1}"
                cell = None
                # XXX: a better data structure to avoid traverse?
                for d in data:
                    if d["class"] == str(classn) and d["coord"] == coord and d["index"] == index:
                        cell = d
                        break
                else:
                    raise LookupError(f"Not found:{classn}/{coord}/{index}")
                table.cell(y, xindex).text = index
                tconversion = {
                    "W": "白",
                    "R": "红",
                    "B": "蓝",
                    "O": "橙",
                    "Y": "黄",
                    "K": "收"
                }
                try:
                    table.cell(y, xdata).text = tconversion[cell["color"]]
                    # If no lyrics, color the data cell
                    coloredcell = table.cell(y, xdata)
                    if lyrics:
                        table.cell(y, xdata).text += "," + lyricslst[i]
                        # Else color the index cell
                        coloredcell = table.cell(y, xindex)
                    coloredcell._tc.get_or_add_tcPr().append(
                        parse_xml(r'<w:shd {} w:fill="{}"/>'.format(
                            nsdecls('w'), cell["colorcode"])))
                except KeyError:
                    print(
                        f"Warning: malformed color code `{cell['color']}' at {classn}/{coord}/{index}", file=sys.stderr)
            # Add a empty line after every table and split every four into two pages
            if person % 8 == 7:
                doc.add_page_break()
            else:
                p = doc.add_paragraph(text='')
                p.paragraph_format.line_spacing = Cm(0.1)
    doc.save(outaddr)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "-s", "--save", help="Save the read cache into a JSON file")
    parser.add_argument(
        "-c", "--cache", help="Read the cache into a JSON file")
    parser.add_argument(
        "-o", "--only", help="The class number to generate, comma separated[all]")
    parser.add_argument("-a", "--analyzeonly",
                        help="Do not generate, analyze only", action="store_true")
    parser.add_argument("-l", "--lyrics", help="Lyrics json file")
    # The m.8 is kept for historical reasons
    parser.add_argument("-n", "--numbering", help="The numbering scheme, `m.8' for the eight beats numbering and `int' for continuing numbers",
                        choices=["m.8", "int"], default="int")
    parser.add_argument("pptfile", help="The input pptx file")
    args = parser.parse_args()
    if not args.cache:
        result, length = readppt(args.pptfile, args.numbering)
    else:
        interm = json.load(open(args.cache))
        result, length = interm[0], interm[1]
    if args.save:
        json.dump([result, length], open(args.save, "w"))
    print("Reading done, generating...")
    if not args.analyzeonly:
        gendoc(result, length, args.numbering, args.pptfile +
               ".docx", [int(c) for c in args.only.split(',')] if args.only else None, args.lyrics)


if __name__ == "__main__":
    main()
