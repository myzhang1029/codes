#!/usr/bin/env python3
# -*- coding: UTF-8 -*-
import pptx, docx
from sys import argv
import json

def readppt(address):
    result = []
    ppt = pptx.Presentation(address)
    for i, slide in enumerate(ppt.slides):
        # The slide index: m.n where 0<n<=8 and m increases
        m = i // 8 + 1
        n = i % 8 + 1
        for shape in slide.shapes: 
            if shape.shape_type == 19: 
                table = shape 
                break
        for nr, row in enumerate(table.table.rows):
            if nr == 0 or nr == 27:
                continue # skip title bar
            for nc, cell in enumerate(row.cells):
                if nc == 0 or nc == 41:
                    continue
                #color = "some magic"
                conversion = {
                    "W":"W",
                    "R":"R",
                    "B":"L",
                    "O":"O",
                    "K":"K"
                }
                try:
                    color = conversion[cell.text]
                except KeyError:
                    color=""
                if 0 < nc < 21: # [1,20], read the left title
                    label = 0
                elif 20 < nc < 41: # [21,40], read the right title
                    label = 41
                classn = row.cells[label].text[:-1] # The class number as a string
                coord = row.cells[label].text[-1] # The row number, which is a alphabet
                coord += table.table.rows[0].cells[nc].text # add the column number
                result.append({
                    "color": color,
                    "class": classn,
                    "coord": coord,
                    "index": f"{m}.{n}"
                    })
    return result

def gendoc(data, outaddr):
    doc = docx.Document()
    for i in range(26*40):
        print(classn := str(i // 40 + 1))
        rem = i % 40
        coord = 'A' if rem < 20 else 'B'
        coord += str(rem % 20 + 1)
        table = doc.add_table(rows=9, cols=8)
        table.style = 'Table Grid'
        r = table.rows[0].cells
        r[0].text = "班级"
        r[1].text = classn
        r[2].text = "编码"
        r[3].text = coord
        r[4].text = "姓名"
        for m in range(1,5):
            for n in range(1,9):
                index = f"{m}.{n}"
                cell = None
                if m == 4 and n > 6:
                    break
                # XXX: a better data structure to avoid traverse?
                for d in data:
                    if d["class"] == classn and d["coord"] == coord and d["index"] == index:
                        cell = d
                        break
                if cell is None:
                    raise IndexError(f"Not found:{classn}/{coord}/{index}")
                if m == 4:
                    if n == 5:
                        index = "国旗"
                    if n == 6:
                        index = "校旗"
                table.rows[n].cells[m*2-2].text = index
                table.rows[n].cells[m*2-1].text = cell["color"]
        if i % 5 == 0 and i != 0:
            doc.add_page_break()
    doc.save(outaddr)

def main():
    result = readppt(argv[1])
    print("Reading done, generating...")
    gendoc(result, argv[1] + ".docx")

if __name__ == "__main__":
    main()
